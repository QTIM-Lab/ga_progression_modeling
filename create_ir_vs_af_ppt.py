import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import cv2
import tempfile
import os
from PIL import Image

df = pd.read_csv('data/GA_progression_modelling_data_redone/clean_data_talisa_ga_10312024.csv')

# get images
slides_data = []
for pat, pat_df in df.groupby('PID'):
    for lat, lat_df in pat_df.groupby('Laterality'):
        for date, date_df in lat_df.groupby('ExamDate'):
            df_af = date_df[date_df.Procedure == 'Af']
            df_ir = date_df[date_df['type'] == 'SLOImage']

            for i, row_af in df_af.iterrows():
                for j, row_ir in df_ir.iterrows():
                    slides_data.append([row_af.file_path_coris, row_ir.file_path_coris, row_af.file_path_ga_seg, row_ir.file_path_ga_seg, pat, lat, date])   

def convert_to_supported_format(image_path):
    """
    Converts an image to a supported format (PNG) using a temporary file.
    Returns the path to the temporary converted file.
    """
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
    try:
        with Image.open(image_path) as img:
            img.convert("RGB").save(temp_file.name, format="PNG")
        return temp_file.name
    except Exception as e:
        print(f"Error converting image {image_path}: {e}")
        raise

def draw_contours(image_path, segmentation_path):
    """
    Draws contours on the image using segmentation and returns the path to the temporary file.
    """
    image = cv2.imread(image_path)
    segmentation = cv2.imread(segmentation_path, cv2.IMREAD_GRAYSCALE)

    # Overlay contours
    contours, _ = cv2.findContours(segmentation, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    overlay = image.copy()
    cv2.drawContours(overlay, contours, -1, (255, 0, 0), thickness=2)

    # Save the output to a temporary file
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
    cv2.imwrite(temp_file.name, overlay)
    return temp_file.name

def create_ppt(slides_data, ppt_output_path):
    prs = Presentation()
    temp_files = []

    try:
        for data in slides_data:
            # Extract information
            image_1_path, image_2_path, image_1_seg, image_2_seg = data[:4]
            patient_id, laterality, exam_date = data[4:]

            # Convert images to supported format
            image_1_path = convert_to_supported_format(image_1_path)
            image_2_path = convert_to_supported_format(image_2_path)
            image_1_seg = convert_to_supported_format(image_1_seg)
            image_2_seg = convert_to_supported_format(image_2_seg)

            # Track temporary files for cleanup
            temp_files.extend([image_1_path, image_2_path, image_1_seg, image_2_seg])

            # Generate overlay images
            image_1_overlay = draw_contours(image_1_path, image_1_seg)
            image_2_overlay = draw_contours(image_2_path, image_2_seg)

            temp_files.extend([image_1_overlay, image_2_overlay])

            # Add a slide
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = f"Patient ID: {patient_id}, Laterality: {laterality}, Exam Date: {exam_date}"

            # Add images to the slide
            top_row_y = Inches(2)
            bottom_row_y = Inches(4.5)

            slide.shapes.add_picture(image_1_path, Inches(1.5), top_row_y, width=Inches(2))
            slide.shapes.add_picture(image_1_seg, Inches(4), top_row_y, width=Inches(2))
            slide.shapes.add_picture(image_1_overlay, Inches(6.5), top_row_y, width=Inches(2))

            slide.shapes.add_picture(image_2_path, Inches(1.5), bottom_row_y, width=Inches(2))
            slide.shapes.add_picture(image_2_seg, Inches(4), bottom_row_y, width=Inches(2))
            slide.shapes.add_picture(image_2_overlay, Inches(6.5), bottom_row_y, width=Inches(2))

        # Save the presentation
        prs.save(ppt_output_path)

    finally:
        # Clean up temporary files
        for temp_file in temp_files:
            try:
                os.remove(temp_file)
            except OSError as e:
                print(f"Error deleting temporary file {temp_file}: {e}")

create_ppt(slides_data, "patient_slides.pptx")
