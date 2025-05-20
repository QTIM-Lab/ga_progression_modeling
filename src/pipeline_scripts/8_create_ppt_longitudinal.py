import os, sys, argparse
import random
import json
import pandas as pd
import torch
from PIL import Image
import numpy as np
import torch.nn.functional as F
import cv2
from moviepy.editor import ImageSequenceClip
import matplotlib.pyplot as plt
import matplotlib.colors as mcolors
from pptx import Presentation
from pptx.util import Inches, Pt
from math import isnan
from datetime import datetime
from torchvision.transforms import Resize, Grayscale, ToTensor
from multiprocessing import Pool, cpu_count
from functools import partial
import tempfile
import shutil
from tqdm import tqdm
import logging

# Configure logging to be minimal
logging.basicConfig(level=logging.WARNING)

def log_progress(task_name, progress, total):
    message = json.dumps({
        "task": task_name,
        "progress": progress,
        "total": total
    })
    print(message)
    sys.stdout.flush()

# def is_valid_date(date_string):
#     formats = ['%m/%d/%Y', '%m-%d-%Y', '%Y-%m-%d', '%Y/%m/%d']
    
#     for date_format in formats:
#         try:
#             # Attempt to parse the date string using the current format
#             datetime.strptime(date_string, date_format)
#             return True
#         except ValueError:
#             continue  # If it fails, try the next format
    
#     return False

def is_date_format(val):
    import re
    return bool(re.match(r"\d{2}-\d{2}-\d{4}", str(val)))

def load_image(path, size=(256, 256), mode='rgb'):
    x = Image.open(path.replace('~', '/home/veturiy'))
    x = Resize(size)(x)
    x = Grayscale()(x) if mode == 'gray' else x
    x = ToTensor()(x)
    return x.unsqueeze(0)

def apply_registration(image_tensor, seg_tensor, grid_path):

    # handle unregistered images
    if isinstance(grid_path, float) and isnan(grid_path):
        # return torch.zeros_like(image_tensor).squeeze(0), torch.zeros_like(seg_tensor).squeeze(0), np.eye(3)
        return image_tensor.squeeze(0), seg_tensor.squeeze(0), np.eye(3)

    # Load the sampling grid
    try:
        params, grid = torch.load(grid_path, weights_only=True)
    
        # Apply the sampling grid
        registered_image = F.grid_sample(image_tensor, grid, align_corners=True)
        registered_seg = F.grid_sample(seg_tensor, grid, align_corners=True)

        # get the affine component
        affine = params[0, -3:, :].T.numpy() # (2, 3)
        affine = np.concatenate([affine[:, 1:], affine[:, :1]], axis=1)
        affine = np.concatenate([affine, np.array([[0, 0, 1]])], axis=0)
        affine = np.linalg.inv(affine)

    except:
        params = torch.load(grid_path, weights_only=True)

        registered_image = torch.permute(image_tensor.squeeze(0), (1, 2, 0)).numpy() # (h, w, c)
        registered_seg = torch.permute(seg_tensor.squeeze(0), (1, 2, 0)).numpy() # (h, w, c)
        
        affine = params.numpy().squeeze(0)
        registered_image = cv2.warpAffine(registered_image, affine[:2, :], (registered_image.shape[0], registered_image.shape[1]))
        registered_seg = cv2.warpAffine(registered_seg, affine[:2, :], (registered_seg.shape[0], registered_seg.shape[1]))

        if registered_image.ndim == 2: # adding extra dim for grayscale warp
            registered_image = registered_image[:, :, None]
        registered_image = torch.tensor(registered_image).permute(2, 0, 1).unsqueeze(0)

        if registered_seg.ndim == 2: # adding extra dim for grayscale warp
            registered_seg = registered_seg[:, :, None]
        registered_seg = torch.tensor(registered_seg).permute(2, 0, 1).unsqueeze(0)
    
    return registered_image.squeeze(0), registered_seg.squeeze(0), affine

def tensor_to_numpy(tensor):
    return tensor.permute(1, 2, 0).numpy()

def draw_contours(image_tensor, seg_tensor, alpha=1.):
    image_np = tensor_to_numpy(image_tensor)
    seg_np = tensor_to_numpy(seg_tensor)

    image_np = (image_np * 255).astype(np.uint8)
    seg_np = (seg_np * 255).astype(np.uint8)

    if image_np.ndim == 2:
        image_np = cv2.cvtColor(image_np, cv2.COLOR_GRAY2BGR)
    
    contours, _ = cv2.findContours(seg_np, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    image_with_contours = cv2.drawContours(image_np.copy(), contours, -1, (0, 255, 0), 1)

    # Blend the overlay with the original image
    cv2.addWeighted(image_with_contours, alpha, image_np, 1 - alpha, 0, image_with_contours)
    
    return torch.tensor(image_with_contours).permute(2, 0, 1)

def visualize_topographical_map(tensor, segmentations, datetimes, cbarlimit=None):
    """
    Visualizes a topographical map with segmentation contours on the baseline image.

    Args:
        tensor (torch.Tensor): A 3D tensor of shape (3, 256, 256) representing the baseline image.
        segmentations (list of torch.Tensor): A list of 2D tensors representing binary segmentation masks.
        datetimes (list of datetime): A list of datetime objects corresponding to each segmentation.
    """
    # Convert datetimes to relative timepoints in years
    baseline_time = datetimes[0]
    if isinstance(baseline_time, datetime):
        relative_timepoints = [(dt - baseline_time).days / 365.25 for dt in datetimes]
    else:
        relative_timepoints = datetimes

    # Convert tensor to a numpy array for visualization
    image_np = tensor_to_numpy(tensor)
    if image_np.ndim == 2:
        image_np = cv2.cvtColor(image_np, cv2.COLOR_GRAY2BGR)
    image_np = (image_np * 255).astype(np.uint8)

    # Set up color map max(relative_timepoints)
    cbarlimit = max(relative_timepoints) if cbarlimit is None else cbarlimit
    norm = mcolors.Normalize(vmin=0, vmax=cbarlimit)
    cmap = plt.cm.viridis

    # Initialize an RGB image for overlaying contours
    contour_overlay = image_np.copy()

    for seg, timepoint in zip(segmentations, relative_timepoints):

        # Convert segmentation tensor to numpy
        seg_np = tensor_to_numpy(seg)
        seg_np = (seg_np * 255).astype(np.uint8)

        # Find contours using OpenCV
        contours, _ = cv2.findContours(seg_np, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

        # Define contour color using the colormap
        contour_color = (np.array(cmap(norm(timepoint))[:3]) * 255).astype(np.uint8).tolist()
        
        cv2.drawContours(contour_overlay, contours, -1, contour_color, 1)

    # Overlay the contours on the baseline image
    overlaid_image = contour_overlay

    # Plot the resulting image
    fig, ax = plt.subplots()
    # cv2.cvtColor(overlaid_image, cv2.COLOR_BGR2RGB)
    ax.imshow(overlaid_image)
    ax.axis('off')

    # Add color bar to the axis
    sm = plt.cm.ScalarMappable(cmap=cmap, norm=norm)
    fig.colorbar(sm, ax=ax, label='Relative Time (years from baseline)')
    ax.set_title('Topographical Map with Segmentation Contours')

def overlap(img1, img2, method='dice'):
    intersection = (img1 * img2).sum()
    min_area = torch.min(img1.sum(), img2.sum())
    return intersection / (min_area + 1e-10)

def deviation(aff):
    return np.linalg.norm(aff[:2, :2] - np.eye(2), 'fro')

def create_videos_and_plots(df, save_to, video_wcontours_name, video_wocontours_name, plot_name, thumbnails_wcontours_name, thumbnails_wocontours_name, baseline_wcontours_name, limits, config):

    # Initialize lists for frames and areas
    frames_with_contours = []
    frames_without_contours = []
    areas = []
    timepoints = []
    longitudinal_segs = []

    for i, row in df.iterrows():
        image_tensor = load_image(row[config.image_col])
        seg_tensor = load_image(row[config.ga_col])
        if i > 0:
            registered_image, registered_seg, aff = apply_registration(image_tensor, seg_tensor, row.params)
            registered_seg = (registered_seg > 0.5).int()
        else:
            registered_image, registered_seg = image_tensor.squeeze(0), seg_tensor.squeeze(0)
            registered_seg = (registered_seg > 0.5).int()
            aff = np.eye(2)

        # remove cases with bad overlap with the baseline image
        # comment if using overlap metric
        # if overlap(registered_seg, baseline_seg) > 0.5:
        # remove cases with drastic transformations in the affine matrix
        
        # filtering regs by affine matrix values
        if config.filter_regs == True:
            if deviation(aff) < 0.1:
                image_with_contours = draw_contours(registered_image, registered_seg)
                longitudinal_segs.append(registered_seg)
                image_without_contours = (registered_image * 255.).int()
        
                # save registered images
                frames_with_contours.append(tensor_to_numpy(image_with_contours))
                frames_without_contours.append(tensor_to_numpy(image_without_contours))

        # not filtering regs
        else:
            image_with_contours = draw_contours(registered_image, registered_seg)
            longitudinal_segs.append(registered_seg)
            image_without_contours = (registered_image * 255.).int()
        
            # save registered images
            frames_with_contours.append(tensor_to_numpy(image_with_contours))
            frames_without_contours.append(tensor_to_numpy(image_without_contours))

        # save area
        areas.append(row[config.area_col])
        timepoints.append(row[config.date_col])

    if len(frames_with_contours) > 0:
        # save thumbnails for videos
        thumbnail_with_contours = frames_with_contours[0].astype(np.uint8)
        thumbnail_without_contours = frames_without_contours[0].astype(np.uint8)
        Image.fromarray(thumbnail_with_contours).save(os.path.join(save_to, thumbnails_wcontours_name))
        Image.fromarray(thumbnail_without_contours).save(os.path.join(save_to, thumbnails_wocontours_name))

        # create videos
        clip_with_contours = ImageSequenceClip(frames_with_contours, fps=3)
        clip_without_contours = ImageSequenceClip(frames_without_contours, fps=3)

        # create buffers and save
        clip_with_contours.write_videofile(os.path.join(save_to, video_wcontours_name), verbose=False, logger=None)
        clip_without_contours.write_videofile(os.path.join(save_to, video_wocontours_name), verbose=False, logger=None)

    # Create a figure with subplots that share the same x-axis
    min_date, max_date, min_area, max_area = limits
    if isinstance(min_date, datetime):
        cbarrange = int((max_date - min_date).days/365.25)+1
    else:
        cbarrange = int(max_date - min_date)+1
    fig, axs = plt.subplots(1, 1, figsize=(7, 8), sharex=True)

    # area vs time
    axs.plot(timepoints, areas, 'ro-')
    axs.legend(['AI-computed Area', 'Manual'])
    axs.set_ylabel('Area (mm$^2$)')  # LaTeX formatting for mm²
    axs.set_xlim(min_date, max_date)
    axs.set_ylim(0, max_area)

    # Rotate x-axis tick labels vertically
    plt.xticks(rotation=90)

    # Save the figure
    plt.savefig(os.path.join(save_to, plot_name))
    plt.close()

    # visualize topological map of GA
    baseline_image = load_image(df.iloc[0][config.image_col]).squeeze(0)
    visualize_topographical_map(baseline_image, longitudinal_segs, timepoints, cbarlimit=cbarrange)
    plt.savefig(os.path.join(save_to, baseline_wcontours_name), dpi=300, bbox_inches='tight')
    plt.close()
    
    return None

def prepare_presentation_slide(slide, slide_title, video_with_contours_path, video_without_contours_path, thumbnail_with_contours_path, thumbnail_without_contours_path, plot_path, config):
    title = slide.shapes.title
    title.text = slide_title
    # title.text_frame.paragraphs[0].font.size = Pt(40)
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(30)

    # Add videos
    video_width = Inches(2.5)
    video_height = Inches(2.5)
    left = Inches(1.0)
    top = Inches(1.75)

    # Add video with contours
    if os.path.exists(video_with_contours_path):
        slide.shapes.add_movie(video_with_contours_path, left, top, width=video_width, height=video_height, poster_frame_image=thumbnail_with_contours_path)

    # Add video without contours
    if os.path.exists(video_without_contours_path):
        slide.shapes.add_movie(video_without_contours_path, left, top + Inches(0.5) + Inches(2.5), width=video_width, height=video_height, poster_frame_image=thumbnail_without_contours_path)

    # Add plot
    plot_width = Inches(4.91)
    plot_height = Inches(5.67)
    # bypass plot path with gompertz
    mrn, lat, _ = os.path.basename(plot_path).rsplit('_', 2)
    if config.gompertz_path is not None:
        if os.path.exists(os.path.join(config.gompertz_path, f'gompertz-fit-{mrn}_{lat}.png')):
            plot_path = os.path.join(config.gompertz_path, f'gompertz-fit-{mrn}_{lat}.png')

    if os.path.exists(plot_path):
        slide.shapes.add_picture(plot_path, Inches(4.5), Inches(1.5), width=plot_width, height=plot_height)

def process_patient_data(mrn_lat_tuple, df, total, limits, config, mrn_mapping):
    """Process a single patient's data and return slide information.
    
    Args:
        mrn_lat_tuple (tuple): Tuple of (mrn, lat) to process
        df (pd.DataFrame): Input dataframe
        limits (tuple): Data limits for plotting
        config: Configuration object
        mrn_mapping (dict): MRN mapping dictionary
    """

    # get patient data
    idx, (mrn, lat), total = mrn_lat_tuple
    lat_df = df[(df[config.patient_col] == mrn) & (df[config.laterality_col] == lat)].copy()

    # log progress
    log_progress('Creating growth curves and videos', idx+1, total)

    # convert to datetime type and sort
    # lat_df[config.date_col] = pd.to_datetime(lat_df[config.date_col])
    lat_df = lat_df.sort_values(by=config.date_col)
    
    # Create temporary directory for this patient's files
    with tempfile.TemporaryDirectory() as temp_dir:

        # Modified paths to use temporary directory
        temp_paths = {
            'videos_wcontours': os.path.join(temp_dir, f'{mrn}_{lat}_wcontours.mp4'),
            'videos_wocontours': os.path.join(temp_dir, f'{mrn}_{lat}_wocontours.mp4'),
            'plots': os.path.join(temp_dir, f'{mrn}_{lat}_plot.png'),
            'thumbnails_wcontours': os.path.join(temp_dir, f'{mrn}_{lat}_wcontours.png'),
            'thumbnails_wocontours': os.path.join(temp_dir, f'{mrn}_{lat}_wocontours.png'),
            'baseline_wcontours': os.path.join(temp_dir, f'{mrn}_{lat}_baseline_wcontours.png')
        }
        # Create videos and plots in temporary directory
        create_videos_and_plots(
            lat_df.reset_index(drop=True),
            save_to=temp_dir,
            video_wcontours_name=f'{mrn}_{lat}_wcontours.mp4',
            video_wocontours_name=f'{mrn}_{lat}_wocontours.mp4',
            plot_name=f'{mrn}_{lat}_plot.png',
            thumbnails_wcontours_name=f'{mrn}_{lat}_wcontours.png',
            thumbnails_wocontours_name=f'{mrn}_{lat}_wocontours.png',
            baseline_wcontours_name=f'{mrn}_{lat}_baseline_wcontours.png',
            limits=limits,
            config=config
        )

        # Create required directories in final destination
        for dir_type in ['videos_wcontours', 'videos_wocontours', 'plots', 'thumbnails_wcontours', 'thumbnails_wocontours', 'baseline_wcontours']:
            os.makedirs(os.path.join(config.output_folder, 'metadata', dir_type), exist_ok=True)

        # Copy files to final destination with debug logging
        for key, temp_path in temp_paths.items():
            if os.path.exists(temp_path):
                final_path = os.path.join(config.output_folder, 'metadata', key, os.path.basename(temp_path))
                try:
                    shutil.copy2(temp_path, final_path)
                except Exception as e:
                    logging.warning(f"Failed to copy {os.path.basename(temp_path)}: {str(e)}")

    return {
        'mrn': mrn,
        'lat': lat,
        'mrn_mapped': mrn_mapping[mrn]
    }

def process_data(df, args):

    # check if valid date
    if is_date_format(df[args.date_col].min()):
        df[args.date_col] = pd.to_datetime(df[args.date_col])
    else:
        df[args.date_col] = df[args.date_col].astype(int)
    
    # Calculate data limits (for plotting)
    limits = (
        datetime(df[args.date_col].min().year - 1, 1, 1) if isinstance(df[args.date_col].min(), datetime) else df[args.date_col].min() - 1,
        datetime(df[args.date_col].max().year + 1, 1, 1) if isinstance(df[args.date_col].max(), datetime) else df[args.date_col].max() + 1,
        df[args.area_col].min() - 10,
        df[args.area_col].max() + 10
    )
    
    # Create MRN mapping (for deidentified data)
    if args.deidentify:
        mrn_mapping = {int(mrn): f'patient{i+1}' for i, mrn in enumerate(df[args.patient_col].unique())}
        with open(os.path.join(args.output_folder, 'metadata', 'mrn_mapping.json'), 'w') as f:
            json.dump(mrn_mapping, f)
    else:
        mrn_mapping = {mrn: mrn for mrn in df[args.patient_col].unique()}
    
    # Get list of all patient-laterality combinations
    groups = df.groupby([args.patient_col, args.laterality_col])
    mrn_lat_pairs = [(i, group, len(groups)) for i, (group, _) in enumerate(groups)]
    
    # # Process patients in parallel
    process_func = partial(process_patient_data, df=df, total=len(mrn_lat_pairs), limits=limits, config=args, mrn_mapping=mrn_mapping)
    with Pool(processes=10) as pool:
        results = list(pool.imap(process_func, mrn_lat_pairs))

    # Filter out None results and create presentation
    results = [r for r in results if r is not None]

    # results = []
    # for pair in mrn_lat_pairs:
    #     results.append(process_func(pair))
    
    return results

def make_presentation(results, args):
    # Create PowerPoint presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    
    for i, result in enumerate(results):
        log_progress('Generating PPT-2 (Longitudinal)', i+1, len(results))

        # Add main slide
        slide = prs.slides.add_slide(slide_layout)
        prepare_presentation_slide(
            slide,
            f'MRN: {result["mrn_mapped"]} | Eye: {result["lat"]}',
            os.path.join(args.output_folder, 'metadata', 'videos_wcontours', f'{result["mrn"]}_{result["lat"]}_wcontours.mp4'),
            os.path.join(args.output_folder, 'metadata', 'videos_wocontours', f'{result["mrn"]}_{result["lat"]}_wocontours.mp4'),
            os.path.join(args.output_folder, 'metadata', 'thumbnails_wcontours', f'{result["mrn"]}_{result["lat"]}_wcontours.png'),
            os.path.join(args.output_folder, 'metadata', 'thumbnails_wocontours', f'{result["mrn"]}_{result["lat"]}_wocontours.png'),
            os.path.join(args.output_folder, 'metadata', 'plots', f'{result["mrn"]}_{result["lat"]}_plot.png'),
            args
        )
        
        # Add topological view slide
        slide = prs.slides.add_slide(slide_layout)
        image_path = os.path.join(args.output_folder, 'metadata', 'baseline_wcontours', f'{result["mrn"]}_{result["lat"]}_baseline_wcontours.png')
        
        if os.path.exists(image_path):
            image = Image.open(image_path)
            
            # Calculate image positioning and sizing
            slide_width = Inches(10)
            slide_height = Inches(7.5)
            image_aspect_ratio = image.width / image.height
            slide_aspect_ratio = slide_width / slide_height
            
            scale_factor = (slide_width / image.width 
                          if image_aspect_ratio > slide_aspect_ratio 
                          else slide_height / image.height)
            
            new_width = image.width * scale_factor
            new_height = image.height * scale_factor
            left = (slide_width - new_width) / 2
            top = (slide_height - new_height) / 2
            
            slide.shapes.add_picture(image_path, left, top, width=new_width, height=new_height)
    
    # Save the presentation
    prs.save(args.save_as)
    return

def main(args):

    os.environ['IMAGEIO_FFMPEG_EXE'] = 'ffmpeg'

    # make directories
    os.makedirs(os.path.join(args.output_folder, 'metadata'), exist_ok=True)
    
    # load data
    df = pd.read_csv(args.csv)

    # process data
    results = process_data(df, args)

    # make presentation
    make_presentation(results, args)

if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('--csv', default='results/test/pipeline_files/images_4.csv', type=str, help='Name of the results file')
    parser.add_argument('--image_col', default='file_path_coris', type=str, help='Column in results file for image.')
    parser.add_argument('--ga_col', default='file_path_ga_seg', type=str, help='Column in results file for segmentation.')
    parser.add_argument('--patient_col', default='PID', type=str, help='Column in results file for patient ID')
    parser.add_argument('--laterality_col', default='Laterality', type=str, help='Column in results file for eye laterality.')
    parser.add_argument('--date_col', default='ExamDate', type=str, help='Column in results file for exam date. Supported types: datetimes for absolute dates, or int for relative time.')
    parser.add_argument('--area_col', default='mm_area', type=str, help='Column in results file for AI extracted lesion area')
    parser.add_argument('--output_folder', default='results/test/pipeline_files/powerpoint_longitudinal', type=str, help='Path to results folder')
    parser.add_argument('--save_as', default='results/test/pipeline_files/powerpoint_longitudinal/powerpoint_longitudinal.pptx', help='Filename to save ppt as')
    parser.add_argument('--deidentify', action='store_true', help='Removes PHI from presentation')
    parser.add_argument('--gompertz_path', default=None, help='Path to gompertz folder.')
    parser.add_argument('--filter_regs', action='store_true', help='Filters out bad registrations by analysing affine matrix')
    args = parser.parse_args()
    main(args)