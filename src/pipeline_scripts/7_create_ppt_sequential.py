import os, sys, argparse
import json
import numpy as np
import pandas as pd
import torch
import cv2
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from math import isnan

DPI = 96  # Standard screen DPI

def deviation(aff):
    return np.linalg.norm(aff[:2, :2] - np.eye(2), 'fro')

def apply_registeration(image_tensor, seg_tensor, grid_path):

    # handle unregistered images
    if isinstance(grid_path, float) and isnan(grid_path):
        return image_tensor, seg_tensor, np.eye(3)

    # resize to (256, 256)
    original_size = image_tensor.shape # (h, w, c)
    image_tensor = cv2.resize(image_tensor, (256, 256)) # (256, 256, c)
    seg_tensor = cv2.resize(seg_tensor, (256, 256)) # (256, 256, c)

    # Load the sampling grid
    try:
        params, grid = torch.load(grid_path, weights_only=True)

        # convert to pytorch tensor
        image_tensor = torch.tensor(image_tensor).permute(2, 0, 1).unsqueeze(0) # (1, c, 256, 256)
        seg_tensor = torch.tensor(seg_tensor).permute(2, 0, 1).unsqueeze(0) # (1, c, 256, 256)
    
        # Apply the sampling grid
        registered_image = F.grid_sample(image_tensor, grid, align_corners=True).squeeze(0).permute(1, 2, 0) # (256, 256, c)
        registered_seg = F.grid_sample(seg_tensor, grid, align_corners=True).squeeze(0).permute(1, 2, 0) # (256, 256, c)

        # get the affine component
        affine = params[0, -3:, :].T.numpy() # (2, 3)
        affine = np.concatenate([affine[:, 1:], affine[:, :1]], axis=1)
        affine = np.concatenate([affine, np.array([[0, 0, 1]])], axis=0)
        affine = np.linalg.inv(affine)

    except:
        params = torch.load(grid_path, weights_only=True)

        # apply affine transform
        affine = params.numpy().squeeze(0)
        registered_image = cv2.warpAffine(image_tensor, affine[:2, :], (image_tensor.shape[0], image_tensor.shape[1]))
        registered_seg = cv2.warpAffine(seg_tensor, affine[:2, :], (seg_tensor.shape[0], seg_tensor.shape[1]))

        if registered_image.ndim == 2: # adding extra dim for grayscale warp
            registered_image = registered_image[:, :, None]

        if registered_seg.ndim == 2: # adding extra dim for grayscale warp
            registered_seg = registered_seg[:, :, None]

    # resize back to original resolution
    registered_image = cv2.resize(registered_image, original_size[:2]) # (h, w, c)
    registered_seg = cv2.resize(registered_seg, original_size[:2]) # (h, w, c)
    
    return registered_image, registered_seg, affine

def log_progress(task_name, progress, total):
    message = json.dumps({
        "task": task_name,
        "progress": progress,
        "total": total
    })
    print(message)
    sys.stdout.flush()

def resize_image_to_fit(image, max_width, max_height):
    """
    Resizes an image to fit within a bounding box (max_width x max_height) while maintaining aspect ratio.
    """
    h, w = image.shape[:2]
    scale_factor = min(max_width / w, max_height / h)
    new_width = int(w * scale_factor)
    new_height = int(h * scale_factor)
    resized_image = cv2.resize(image, (new_width, new_height))
    return resized_image

def process_images(original_img_path, seg_img_path, slide_width, slide_height, params=None):
    """
    Prepares the original image, segmentation, and overlay after resizing to fit slide dimensions.
    """
    # Load images
    original_img = np.array(Image.open(original_img_path)) # (h, w, c)
    seg_img = np.array(Image.open(seg_img_path).resize((original_img.shape[1], original_img.shape[0]))) # (h, w, c)

    # Register
    if params:
        original_img_reg, seg_img_reg, affine = apply_registeration(original_img, seg_img, params)
        if deviation(affine) < 0.1:
            original_img = original_img_reg
            seg_img = seg_img_reg

    # Resize images to fit half of slide height
    max_width = slide_width // 3  # Divide slide width by 3 for three images
    max_height = slide_height // 2  # Allow some space for the title
    original_resized = resize_image_to_fit(original_img, max_width, max_height)
    seg_resized = resize_image_to_fit(seg_img, max_width, max_height)

    # Binarize the segmentation
    _, seg_binarized = cv2.threshold(seg_resized, 1, 255, cv2.THRESH_BINARY)

    # Create overlay with contours
    overlay = original_resized.copy()
    contours, _ = cv2.findContours(seg_binarized, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    cv2.drawContours(overlay, contours, -1, (0, 255, 0), 1)  # Blue contours

    return original_resized, seg_binarized, overlay

def main(args):

    # make a metadata folder
    metadata_folder = os.path.join(args.output_folder, 'metadata')
    os.makedirs(metadata_folder, exist_ok=True)

    # Initialize PowerPoint presentation
    presentation = Presentation()

    # Slide dimensions in pixels (convert from inches)
    slide_width_pixels = int(presentation.slide_width.inches * DPI)
    slide_height_pixels = int(presentation.slide_height.inches * DPI)

    # Load CSV data
    df = pd.read_csv(args.csv)

    for idx, row in df.iterrows():
        log_progress('Generating PPT-1 (Indivitual segs)', idx+1, len(df))

        file_path = row[args.image_col]
        seg_path = row[args.ga_col]
        pid = row[args.patient_col]
        laterality = row[args.laterality_col]
        exam_date = row[args.date_col]
        mm_area = row[args.area_col]
        params = row.params if args.register else None

        # Process images
        original_resized, seg_binarized, overlay = process_images(
            file_path, seg_path, slide_width_pixels, slide_height_pixels, params
        )

        # Save temporary files
        cv2.imwrite(os.path.join(metadata_folder, "original_temp.png"), original_resized)
        cv2.imwrite(os.path.join(metadata_folder, "seg_temp.png"), seg_binarized)
        cv2.imwrite(os.path.join(metadata_folder, "overlay_temp.png"), overlay)

        # Add slide and title
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        # slide.shapes.title.text = f"PID: {pid}, Laterality: {laterality}, ExamDate: {exam_date}, mm² Area: {mm_area}"

        title = slide.shapes.title
        title.text = f"PID: {pid}, Laterality: {laterality}, \nExamDate: {exam_date}, Area: {mm_area:.3f} mm² "
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(30)

        # Add images to slide
        left_margin = Inches(0.25)
        top_margin = Inches(2.5)
        image_spacing = Inches(0.25)
        image_width = Inches(3)  # Set consistent width for all images

        slide.shapes.add_picture(os.path.join(metadata_folder, "original_temp.png"), left_margin, top_margin, width=image_width)
        slide.shapes.add_picture(os.path.join(metadata_folder, "seg_temp.png"), left_margin + image_width + image_spacing, top_margin, width=image_width)
        slide.shapes.add_picture(os.path.join(metadata_folder, "overlay_temp.png"), left_margin + 2 * (image_width + image_spacing), top_margin, width=image_width)

    # Save the presentation
    presentation.save(args.save_as)
    print("Updated PowerPoint presentation created successfully!")

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description="Generates sequential presentation from segmentations")
    parser.add_argument('--csv', default='results/test/images_processed.csv')
    parser.add_argument('--image_col', default='file_path_coris', type=str, help='Column in results file for image.')
    parser.add_argument('--ga_col', default='file_path_ga_seg', type=str, help='Column in results file for GA segmentation.')
    parser.add_argument('--patient_col', default='PID', type=str, help='Column in results file for patient ID')
    parser.add_argument('--laterality_col', default='Laterality', type=str, help='Column in results file for eye laterality.')
    parser.add_argument('--date_col', default='ExamDate', type=str, help='Column in results file for exam date. Supported types: datetimes for absolute dates, or int for relative time.')
    parser.add_argument('--area_col', default='mm_area', type=str, help='Column in results file for AI extracted lesion area')
    parser.add_argument('--output_folder', default='results/test/pipeline_files/powerpoint_sequential', type=str, help='Path to results folder')
    parser.add_argument('--register', action='store_true', help='Register images before putting onto the slides')
    parser.add_argument('--save_as', default='results/test/pipeline_files/powerpoint_sequential/sequential.pptx', type=str)
    args = parser.parse_args()
    main(args)