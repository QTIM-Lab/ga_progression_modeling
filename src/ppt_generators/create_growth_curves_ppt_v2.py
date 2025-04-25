import os
import sys
import random
import json
import argparse
import pandas as pd
import torch
from PIL import Image
import numpy as np
import torch.nn.functional as F
import cv2
from moviepy.editor import ImageSequenceClip
import matplotlib.pyplot as plt
import matplotlib.colors as mcolors
from pptx import Presentation
from pptx.util import Inches, Pt
from math import isnan
from datetime import datetime
from torchvision.transforms import Resize, Grayscale, ToTensor
from multiprocessing import Pool, cpu_count
from functools import partial
import tempfile
import shutil
from tqdm import tqdm
import logging

# Configure logging to be minimal
logging.basicConfig(level=logging.WARNING)

def parse_args():
    parser = argparse.ArgumentParser()
    parser.add_argument('--results-folder', type=str, help='Path to results folder')
    parser.add_argument('--results-file', default='results.csv', type=str, help='Name of the results file')
    parser.add_argument('--img-col', type=str, help='Column in results file for image.')
    parser.add_argument('--seg-col', type=str, help='Column in results file for segmentation.')
    parser.add_argument('--patient-col', type=str, help='Column in results file for patient ID')
    parser.add_argument('--laterality-col', type=str, help='Column in results file for eye laterality.')
    parser.add_argument('--date-col', type=str, help='Column in results file for exam date. Supported types: datetimes for absolute dates, or int for relative time.')
    parser.add_argument('--area-manual-col', type=str, default=None, help='Column in results file for manually extracted lesion area')
    parser.add_argument('--area-ai-col', type=str, default=None, help='Column in results file for AI extracted lesion area')
    parser.add_argument('--perimeter-ai-col', type=str, default=None, help='Column in results file for AI extracted lesion perimeter')
    parser.add_argument('--n-foci-ai-col', type=str, default=None, help='Column in results file for AI extracted number of lesion foci')
    parser.add_argument('--multiple-visits', action='store_true', help='Only save slides for patients with multiple visits.')
    parser.add_argument('--specific-pat', type=str, default=None, help='Path to txt file with specific patients to get slides on.')
    parser.add_argument('--ppt-folder', default='powerpoint', help='Folder to save ppt to.')
    parser.add_argument('--ppt-file', default='segmentation_analysis.pptx', help='Filename to save ppt as')
    parser.add_argument('--deidentify', action='store_true', help='Removes PHI from presentation')
    parser.add_argument('--gompertz', default=None, help='Pull plots from gompertz folder.')
    parser.add_argument('--filter_regs', default=False, help='Filters out bad registrations by analysing affine matrix')
    args = parser.parse_args()

    args.filter_regs = False if args.filter_regs == 'false' else True
    return args

def is_valid_date(date_string):
    formats = ['%m/%d/%Y', '%m-%d-%Y', '%Y-%m-%d', '%Y/%m/%d']
    
    for date_format in formats:
        try:
            # Attempt to parse the date string using the current format
            datetime.strptime(date_string, date_format)
            return True
        except ValueError:
            continue  # If it fails, try the next format
    
    return False

def load_image(path, size=(256, 256), mode='rgb'):
    x = Image.open(path.replace('~', '/home/veturiy'))
    x = Resize(size)(x)
    x = Grayscale()(x) if mode == 'gray' else x
    x = ToTensor()(x)
    return x.unsqueeze(0)

def apply_registration(image_tensor, seg_tensor, grid_path):

    # handle unregistered images
    if isinstance(grid_path, float) and isnan(grid_path):
        # return torch.zeros_like(image_tensor).squeeze(0), torch.zeros_like(seg_tensor).squeeze(0), np.eye(3)
        return image_tensor.squeeze(0), seg_tensor.squeeze(0), np.eye(3)

    # Load the sampling grid
    try:
        params, grid = torch.load(grid_path, weights_only=True)
    
        # Apply the sampling grid
        registered_image = F.grid_sample(image_tensor, grid, align_corners=True)
        registered_seg = F.grid_sample(seg_tensor, grid, align_corners=True)

        # get the affine component
        affine = params[0, -3:, :].T.numpy() # (2, 3)
        affine = np.concatenate([affine[:, 1:], affine[:, :1]], axis=1)
        affine = np.concatenate([affine, np.array([[0, 0, 1]])], axis=0)
        affine = np.linalg.inv(affine)

    except:
        params = torch.load(grid_path, weights_only=True)

        registered_image = torch.permute(image_tensor.squeeze(0), (1, 2, 0)).numpy() # (h, w, c)
        registered_seg = torch.permute(seg_tensor.squeeze(0), (1, 2, 0)).numpy() # (h, w, c)
        
        affine = params.numpy().squeeze(0)
        registered_image = cv2.warpAffine(registered_image, affine[:2, :], (registered_image.shape[0], registered_image.shape[1]))
        registered_seg = cv2.warpAffine(registered_seg, affine[:2, :], (registered_seg.shape[0], registered_seg.shape[1]))

        if registered_image.ndim == 2: # adding extra dim for grayscale warp
            registered_image = registered_image[:, :, None]
        registered_image = torch.tensor(registered_image).permute(2, 0, 1).unsqueeze(0)

        if registered_seg.ndim == 2: # adding extra dim for grayscale warp
            registered_seg = registered_seg[:, :, None]
        registered_seg = torch.tensor(registered_seg).permute(2, 0, 1).unsqueeze(0)
    
    return registered_image.squeeze(0), registered_seg.squeeze(0), affine

def tensor_to_numpy(tensor):
    return tensor.permute(1, 2, 0).numpy()

def draw_contours(image_tensor, seg_tensor, alpha=1.):
    image_np = tensor_to_numpy(image_tensor)
    seg_np = tensor_to_numpy(seg_tensor)

    image_np = (image_np * 255).astype(np.uint8)
    seg_np = (seg_np * 255).astype(np.uint8)

    if image_np.ndim == 2:
        image_np = cv2.cvtColor(image_np, cv2.COLOR_GRAY2BGR)
    
    contours, _ = cv2.findContours(seg_np, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    image_with_contours = cv2.drawContours(image_np.copy(), contours, -1, (0, 255, 0), 1)

    # Blend the overlay with the original image
    cv2.addWeighted(image_with_contours, alpha, image_np, 1 - alpha, 0, image_with_contours)
    
    return torch.tensor(image_with_contours).permute(2, 0, 1)

def visualize_topographical_map(tensor, segmentations, datetimes, cbarlimit=None):
    """
    Visualizes a topographical map with segmentation contours on the baseline image.

    Args:
        tensor (torch.Tensor): A 3D tensor of shape (3, 256, 256) representing the baseline image.
        segmentations (list of torch.Tensor): A list of 2D tensors representing binary segmentation masks.
        datetimes (list of datetime): A list of datetime objects corresponding to each segmentation.
    """
    # Convert datetimes to relative timepoints in years
    baseline_time = datetimes[0]
    if isinstance(baseline_time, datetime):
        relative_timepoints = [(dt - baseline_time).days / 365.25 for dt in datetimes]
    else:
        relative_timepoints = datetimes

    # Convert tensor to a numpy array for visualization
    image_np = tensor_to_numpy(tensor)
    if image_np.ndim == 2:
        image_np = cv2.cvtColor(image_np, cv2.COLOR_GRAY2BGR)
    image_np = (image_np * 255).astype(np.uint8)

    # Set up color map max(relative_timepoints)
    cbarlimit = max(relative_timepoints) if cbarlimit is None else cbarlimit
    norm = mcolors.Normalize(vmin=0, vmax=cbarlimit)
    cmap = plt.cm.viridis

    # Initialize an RGB image for overlaying contours
    contour_overlay = image_np.copy()

    for seg, timepoint in zip(segmentations, relative_timepoints):

        # Convert segmentation tensor to numpy
        seg_np = tensor_to_numpy(seg)
        seg_np = (seg_np * 255).astype(np.uint8)

        # Find contours using OpenCV
        contours, _ = cv2.findContours(seg_np, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

        # Define contour color using the colormap
        contour_color = (np.array(cmap(norm(timepoint))[:3]) * 255).astype(np.uint8).tolist()
        
        cv2.drawContours(contour_overlay, contours, -1, contour_color, 1)

    # Overlay the contours on the baseline image
    overlaid_image = contour_overlay

    # Plot the resulting image
    fig, ax = plt.subplots()
    # cv2.cvtColor(overlaid_image, cv2.COLOR_BGR2RGB)
    ax.imshow(overlaid_image)
    ax.axis('off')

    # Add color bar to the axis
    sm = plt.cm.ScalarMappable(cmap=cmap, norm=norm)
    fig.colorbar(sm, ax=ax, label='Relative Time (years from baseline)')
    ax.set_title('Topographical Map with Segmentation Contours')

def overlap(img1, img2, method='dice'):
    intersection = (img1 * img2).sum()
    min_area = torch.min(img1.sum(), img2.sum())
    return intersection / (min_area + 1e-10)

def deviation(aff):
    return np.linalg.norm(aff[:2, :2] - np.eye(2), 'fro')

def create_videos_and_plots(df, save_to, video_wcontours_name, video_wocontours_name, plot_name, thumbnails_wcontours_name, thumbnails_wocontours_name, baseline_wcontours_name, limits, config):

    # Initialize lists for frames and areas
    frames_with_contours = []
    frames_without_contours = []
    areas_ai = []
    areas_manual = []
    perimeters_ai = []
    n_foci_ai = []
    timepoints = []
    longitudinal_segs = []

    for i, row in df.iterrows():
        image_tensor = load_image(row[config.img_col])
        seg_tensor = load_image(row[config.seg_col])
        if i > 0:
            registered_image, registered_seg, aff = apply_registration(image_tensor, seg_tensor, row.params)
            registered_seg = (registered_seg > 0.5).int()
        else:
            registered_image, registered_seg = image_tensor.squeeze(0), seg_tensor.squeeze(0)
            registered_seg = (registered_seg > 0.5).int()
            # baseline_seg = registered_seg.clone()
            aff = np.eye(2)

        # remove cases with bad overlap with the baseline image
        # comment if using overlap metric
        # if overlap(registered_seg, baseline_seg) > 0.5:
        # remove cases with drastic transformations in the affine matrix
        
        # filtering regs
        if config.filter_regs:
            if deviation(aff) < 0.1:
                image_with_contours = draw_contours(registered_image, registered_seg)
                longitudinal_segs.append(registered_seg)
                image_without_contours = (registered_image * 255.).int()
        
                # save registered images
                frames_with_contours.append(tensor_to_numpy(image_with_contours))
                frames_without_contours.append(tensor_to_numpy(image_without_contours))

        # not filtering regs
        else:
            image_with_contours = draw_contours(registered_image, registered_seg)
            longitudinal_segs.append(registered_seg)
            image_without_contours = (registered_image * 255.).int()
        
            # save registered images
            frames_with_contours.append(tensor_to_numpy(image_with_contours))
            frames_without_contours.append(tensor_to_numpy(image_without_contours))

        # compute AI area
        if config.area_ai_col is not None:
            areas_ai.append(row[config.area_ai_col])
        else:
            seg_np = tensor_to_numpy((seg_tensor > 0.5).int().squeeze(0))
            seg_np = cv2.resize(seg_np.astype(np.uint8), (int(row.xslo), int(row.yslo)))
            area = np.sum(seg_np) * row.scale_x * row.scale_y
            areas_ai.append(area)

        # save manual area
        if config.area_manual_col in row.keys():
            areas_manual.append(row[config.area_manual_col])
        else:
            areas_manual.append(float('nan'))

        # compute perimeter
        if config.perimeter_ai_col is not None:
            perimeters_ai.append(row[config.perimeter_ai_col])
        else:
            perimeters_ai.append(None)

        # compute number of foci
        if config.n_foci_ai_col is not None:
            n_foci_ai.append(row[config.n_foci_ai_col])
        else:
            n_foci_ai.append(None)

        timepoints.append(row[config.date_col])

    if len(frames_with_contours) > 0:
        # save thumbnails for videos
        thumbnail_with_contours = frames_with_contours[0].astype(np.uint8)
        thumbnail_without_contours = frames_without_contours[0].astype(np.uint8)
        Image.fromarray(thumbnail_with_contours).save(os.path.join(save_to, thumbnails_wcontours_name))
        Image.fromarray(thumbnail_without_contours).save(os.path.join(save_to, thumbnails_wocontours_name))

        # create videos
        clip_with_contours = ImageSequenceClip(frames_with_contours, fps=3)
        clip_without_contours = ImageSequenceClip(frames_without_contours, fps=3)

        # create buffers and save
        clip_with_contours.write_videofile(os.path.join(save_to, video_wcontours_name), verbose=False, logger=None)
        clip_without_contours.write_videofile(os.path.join(save_to, video_wocontours_name), verbose=False, logger=None)

    # Create a figure with subplots that share the same x-axis
    min_date, max_date, min_area, max_area, min_peri, max_peri, min_foci, max_foci = limits
    fig, axs = plt.subplots(1, 1, figsize=(7, 8), sharex=True)

    # area vs time
    axs.plot(timepoints, areas_ai, 'ro-')
    axs.plot(timepoints, areas_manual, 'bo-')
    axs.legend(['AI-computed Area', 'Manual'])
    axs.set_ylabel('GA Area (mm$^2$)')  # LaTeX formatting for mm²
    axs.set_xlim(min_date, max_date)
    axs.set_ylim(0, max_area)

    # Rotate x-axis tick labels vertically
    plt.xticks(rotation=90)

    # Save the figure
    plt.savefig(os.path.join(save_to, plot_name))
    plt.close()

    # visualize topological map of GA
    baseline_image = load_image(df.iloc[0][config.img_col]).squeeze(0)
    visualize_topographical_map(baseline_image, longitudinal_segs, timepoints, cbarlimit=int((max_date - min_date).days/365.25)+1)
    plt.savefig(os.path.join(save_to, baseline_wcontours_name), dpi=300, bbox_inches='tight')
    plt.close()
    
    return None

def prepare_presentation_slide(slide, slide_title, video_with_contours_path, video_without_contours_path, thumbnail_with_contours_path, thumbnail_without_contours_path, plot_path, config):
    title = slide.shapes.title
    title.text = slide_title
    # title.text_frame.paragraphs[0].font.size = Pt(40)
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(30)

    # Add videos
    video_width = Inches(2.5)
    video_height = Inches(2.5)
    left = Inches(1.0)
    top = Inches(1.75)

    # Add video with contours
    if os.path.exists(video_with_contours_path):
        slide.shapes.add_movie(video_with_contours_path, left, top, width=video_width, height=video_height, poster_frame_image=thumbnail_with_contours_path)

    # Add video without contours
    if os.path.exists(video_without_contours_path):
        slide.shapes.add_movie(video_without_contours_path, left, top + Inches(0.5) + Inches(2.5), width=video_width, height=video_height, poster_frame_image=thumbnail_without_contours_path)

    # Add plot
    plot_width = Inches(4.91)
    plot_height = Inches(5.67)
    # bypass plot path with gompertz
    mrn, lat, _ = os.path.basename(plot_path).split('_')
    if config.gompertz is not None:
        if os.path.exists(os.path.join(config.results_folder, config.gompertz, f'gompertz_plots/gompertz-fit-{mrn}_{lat}.png')):
            plot_path = os.path.join(config.results_folder, config.gompertz, f'gompertz_plots/gompertz-fit-{mrn}_{lat}.png')

    if os.path.exists(plot_path):
        slide.shapes.add_picture(plot_path, Inches(4.5), Inches(1.5), width=plot_width, height=plot_height)

def process_patient_data(mrn_lat_tuple, df, save_to, limits, config, mrn_mapping):
    """Process a single patient's data and return slide information.
    
    Args:
        mrn_lat_tuple (tuple): Tuple of (mrn, lat) to process
        df (pd.DataFrame): Input dataframe
        save_to (str): Base directory to save files
        limits (tuple): Data limits for plotting
        config: Configuration object
        mrn_mapping (dict): MRN mapping dictionary
    """

    mrn, lat = mrn_lat_tuple
    lat_df = df[(df[config.patient_col] == mrn) & (df[config.laterality_col] == lat)].copy()
    
    # Skip processing if conditions aren't met
    if config.multiple_visits and len(lat_df.ExamDate.unique()) == 1:
        return None
    
    if config.specific_pat:
        with open(os.path.join(config.results_folder, config.specific_pat), 'r') as f:
            specific_patients = f.read().splitlines()
        if f'{mrn}_{lat}' not in specific_patients:
            return None

    lat_df[config.date_col] = pd.to_datetime(lat_df[config.date_col])
    lat_df = lat_df.sort_values(by=config.date_col)
    
    # Get patient metadata
    age_v1 = ('NA' if isinstance(lat_df.iloc[0].dob, float) else 
              int((lat_df.iloc[0][config.date_col] - pd.to_datetime(lat_df.iloc[0].dob)).days / 365.25))
    pr_score = 'NA' if isnan(lat_df.iloc[0].PGS004606) else lat_df.iloc[0].PGS004606

    # Create temporary directory for this patient's files
    with tempfile.TemporaryDirectory() as temp_dir:

        # Modified paths to use temporary directory
        temp_paths = {
            'videos_wcontours': os.path.join(temp_dir, f'{mrn}_{lat}_wcontours.mp4'),
            'videos_wocontours': os.path.join(temp_dir, f'{mrn}_{lat}_wocontours.mp4'),
            'plots': os.path.join(temp_dir, f'{mrn}_{lat}_plot.png'),
            'thumbnails_wcontours': os.path.join(temp_dir, f'{mrn}_{lat}_wcontours.png'),
            'thumbnails_wocontours': os.path.join(temp_dir, f'{mrn}_{lat}_wocontours.png'),
            'baseline_wcontours': os.path.join(temp_dir, f'{mrn}_{lat}_baseline_wcontours.png')
        }
        
        # Create videos and plots in temporary directory
        create_videos_and_plots(
            lat_df.reset_index(drop=True),
            save_to=temp_dir,
            video_wcontours_name=f'{mrn}_{lat}_wcontours.mp4',
            video_wocontours_name=f'{mrn}_{lat}_wocontours.mp4',
            plot_name=f'{mrn}_{lat}_plot.png',
            thumbnails_wcontours_name=f'{mrn}_{lat}_wcontours.png',
            thumbnails_wocontours_name=f'{mrn}_{lat}_wocontours.png',
            baseline_wcontours_name=f'{mrn}_{lat}_baseline_wcontours.png',
            limits=limits,
            config=config
        )

        # Create required directories in final destination
        for dir_type in ['videos_wcontours', 'videos_wocontours', 'plots', 'thumbnails_wcontours', 
                        'thumbnails_wocontours', 'baseline_wcontours']:
            os.makedirs(os.path.join(save_to, 'metadata', dir_type), exist_ok=True)

        # Copy files to final destination with debug logging
        for key, temp_path in temp_paths.items():
            if os.path.exists(temp_path):
                final_path = os.path.join(save_to, 'metadata', key, os.path.basename(temp_path))
                try:
                    shutil.copy2(temp_path, final_path)
                except Exception as e:
                    logging.warning(f"Failed to copy {os.path.basename(temp_path)}: {str(e)}")

    return {
        'mrn': mrn,
        'lat': lat,
        'mrn_mapped': mrn_mapping[mrn],
        'age_v1': age_v1,
        'pr_score': pr_score,
        'in_registry': lat_df.in_registry.iloc[0]
    }

def main():
    # Suppress moviepy's verbose output
    import os
    os.environ['IMAGEIO_FFMPEG_EXE'] = 'ffmpeg'

    args = parse_args()
    
    # Load and prepare data
    file = os.path.join(args.results_folder, args.results_file)
    save_to = os.path.join(args.results_folder, args.ppt_folder)
    
    # Clean directory if exists
    if os.path.exists(save_to):
        shutil.rmtree(save_to)
    
    os.makedirs(os.path.join(save_to, 'metadata'), exist_ok=True)
    
    # Read and process dataframe
    df = pd.read_csv(file)
    if is_valid_date(df[args.date_col].min()):
        df[args.date_col] = pd.to_datetime(df[args.date_col])
    
    # Calculate data limits
    limits = (
        datetime(df[args.date_col].min().year - 1, 1, 1) if isinstance(df[args.date_col].min(), datetime) else df[args.date_col].min() - 1,
        datetime(df[args.date_col].max().year + 1, 1, 1) if isinstance(df[args.date_col].min(), datetime) else df[args.date_col].min() + 1,
        df[args.area_ai_col].min() - 10,
        df[args.area_ai_col].max() + 10,
        df[args.perimeter_ai_col].min() - 10,
        df[args.perimeter_ai_col].max() + 10,
        df[args.n_foci_ai_col].min() - 5,
        df[args.n_foci_ai_col].max() + 5
    )
    
    # Create MRN mapping
    if args.deidentify:
        mrn_mapping = {int(mrn): f'patient{i+1}' for i, mrn in enumerate(df[args.patient_col].unique())}
        with open(os.path.join(save_to, 'metadata/mrn_mapping.json'), 'w') as f:
            json.dump(mrn_mapping, f)
    else:
        mrn_mapping = {mrn: mrn for mrn in df[args.patient_col].unique()}
    
    # Get list of all patient-laterality combinations
    mrn_lat_pairs = [(mrn, lat) for mrn, mrn_df in df.groupby(args.patient_col) 
                     for lat in mrn_df[args.laterality_col].unique()] #[:10]
    
    # Process patients in parallel
    process_func = partial(process_patient_data, 
                         df=df, 
                         save_to=save_to, 
                         limits=limits, 
                         config=args, 
                         mrn_mapping=mrn_mapping)
    
    with Pool(processes=cpu_count()) as pool:
        results = list(tqdm(
            pool.imap(process_func, mrn_lat_pairs),
            total=len(mrn_lat_pairs),
            desc="Processing patients",
            unit="patient"
        ))
    
    # Filter out None results and create presentation
    results = [r for r in results if r is not None]
    
    # Create PowerPoint presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    
    for result in tqdm(results, desc="Creating slides", unit="slide"):
        # Add main slide
        slide = prs.slides.add_slide(slide_layout)
        prepare_presentation_slide(
            slide,
            f'MRN: {result["mrn_mapped"]} | Eye: {result["lat"]} | Age(v1): {result["age_v1"]} | PRS: {result["pr_score"]} | In registry?: {result["in_registry"]}',
            os.path.join(save_to, 'metadata', 'videos_wcontours', f'{result["mrn"]}_{result["lat"]}_wcontours.mp4'),
            os.path.join(save_to, 'metadata', 'videos_wocontours', f'{result["mrn"]}_{result["lat"]}_wocontours.mp4'),
            os.path.join(save_to, 'metadata', 'thumbnails_wcontours', f'{result["mrn"]}_{result["lat"]}_wcontours.png'),
            os.path.join(save_to, 'metadata', 'thumbnails_wocontours', f'{result["mrn"]}_{result["lat"]}_wocontours.png'),
            os.path.join(save_to, 'metadata', 'plots', f'{result["mrn"]}_{result["lat"]}_plot.png'),
            args
        )
        
        # Add topological view slide
        slide = prs.slides.add_slide(slide_layout)
        image_path = os.path.join(save_to, 'metadata', 'baseline_wcontours', f'{result["mrn"]}_{result["lat"]}_baseline_wcontours.png')
        
        if os.path.exists(image_path):
            image = Image.open(image_path)
            
            # Calculate image positioning and sizing
            slide_width = Inches(10)
            slide_height = Inches(7.5)
            image_aspect_ratio = image.width / image.height
            slide_aspect_ratio = slide_width / slide_height
            
            scale_factor = (slide_width / image.width 
                          if image_aspect_ratio > slide_aspect_ratio 
                          else slide_height / image.height)
            
            new_width = image.width * scale_factor
            new_height = image.height * scale_factor
            left = (slide_width - new_width) / 2
            top = (slide_height - new_height) / 2
            
            slide.shapes.add_picture(image_path, left, top, width=new_width, height=new_height)
    
    # Save the presentation
    prs.save(os.path.join(save_to, args.ppt_file))

if __name__ == '__main__':
    main()