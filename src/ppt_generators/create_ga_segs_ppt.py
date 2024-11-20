import pandas as pd
import cv2
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Cm, Pt

DPI = 96  # Standard screen DPI

def resize_image_to_fit(image, max_width, max_height):
    """
    Resizes an image to fit within a bounding box (max_width x max_height) while maintaining aspect ratio.
    """
    h, w = image.shape[:2]
    scale_factor = min(max_width / w, max_height / h)
    new_width = int(w * scale_factor)
    new_height = int(h * scale_factor)
    resized_image = cv2.resize(image, (new_width, new_height))
    return resized_image

def process_images(original_img_path, seg_img_path, slide_width, slide_height):
    """
    Prepares the original image, segmentation, and overlay after resizing to fit slide dimensions.
    """
    # Load images
    original_img = cv2.imread(original_img_path)
    seg_img = cv2.imread(seg_img_path, cv2.IMREAD_GRAYSCALE)

    # Resize images to fit half of slide height
    max_width = slide_width // 3  # Divide slide width by 3 for three images
    max_height = slide_height // 2  # Allow some space for the title
    original_resized = resize_image_to_fit(original_img, max_width, max_height)
    seg_resized = resize_image_to_fit(seg_img, max_width, max_height)

    # Binarize the segmentation
    _, seg_binarized = cv2.threshold(seg_resized, 1, 255, cv2.THRESH_BINARY)

    # Create overlay with contours
    overlay = original_resized.copy()
    contours, _ = cv2.findContours(seg_binarized, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    cv2.drawContours(overlay, contours, -1, (0, 255, 0), 1)  # Blue contours

    return original_resized, seg_binarized, overlay

# Initialize PowerPoint presentation
presentation = Presentation()

# Slide dimensions in pixels (convert from inches)
slide_width_pixels = int(presentation.slide_width.inches * DPI)
slide_height_pixels = int(presentation.slide_height.inches * DPI)

# Load CSV data
csv_path = "results/11132024_talisa_ga/area_comparisons_af.csv"  # Update with the actual path
df = pd.read_csv(csv_path)

for _, row in df.iterrows():
    file_path = row['file_path_coris']
    seg_path = row['file_path_ga_seg']
    pid = row['PID']
    laterality = row['Laterality']
    exam_date = row['ExamDate']
    mm_area = row['mm_area']

    # Process images
    original_resized, seg_binarized, overlay = process_images(
        file_path, seg_path, slide_width_pixels, slide_height_pixels
    )

    # Save temporary files
    cv2.imwrite("original_temp.png", original_resized)
    cv2.imwrite("seg_temp.png", seg_binarized)
    cv2.imwrite("overlay_temp.png", overlay)

    # Add slide and title
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    # slide.shapes.title.text = f"PID: {pid}, Laterality: {laterality}, ExamDate: {exam_date}, mm² Area: {mm_area}"

    title = slide.shapes.title
    title.text = f"PID: {pid}, Laterality: {laterality}, \nExamDate: {exam_date}, Area: {mm_area:.3f} mm² "
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(30)

    # Add images to slide
    left_margin = Inches(0.25)
    top_margin = Inches(2.5)
    image_spacing = Inches(0.25)
    image_width = Inches(3)  # Set consistent width for all images

    slide.shapes.add_picture("original_temp.png", left_margin, top_margin, width=image_width)
    slide.shapes.add_picture("seg_temp.png", left_margin + image_width + image_spacing, top_margin, width=image_width)
    slide.shapes.add_picture("overlay_temp.png", left_margin + 2 * (image_width + image_spacing), top_margin, width=image_width)

# Save the presentation
presentation.save("GA_images.pptx")
print("Updated PowerPoint presentation created successfully!")
