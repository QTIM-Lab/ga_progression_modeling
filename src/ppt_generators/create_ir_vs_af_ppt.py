import os, sys
sys.path.append('/sddata/projects/LightGlue/src/')

import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import cv2
import numpy as np
import tempfile
import torch
from PIL import Image
from tqdm import tqdm
from concurrent.futures import ProcessPoolExecutor
from eyeliner import EyeLinerP
from eyeliner.utils import load_image

device = 'cpu'

# Load EyeLiner API for registration
eyelinerp = EyeLinerP(
    reg='affine', # registration technique to use (tps or affine)
    lambda_tps=1.0, # set lambda value for tps
    image_size=(3, 256, 256), # image dimensions
    device=device
    )

def tensor2numpy(tensor):
    return torch.permute(tensor, (1, 2, 0)).numpy()

def register(image_1_path, image_2_path, image_1_seg, image_2_seg):

    w, h = Image.open(image_1_path).size
    w_slo, h_slo = Image.open(image_2_path).size

    # load each image
    slo_image = load_image(image_2_path, size=(256, 256), mode='rgb').to(device) # (256, 256), [0, 1]
    slo_ga_seg = load_image(image_2_seg, size=(256, 256), mode='rgb').to(device) # (256, 256), [0, 1]

    faf_image = load_image(image_1_path, size=(256, 256), mode='rgb').to(device) # (256, 256), [0, 1]
    faf_ga_seg = load_image(image_1_seg, size=(256, 256), mode='rgb').to(device) # (256, 256), [0, 1]

    # store inputs
    data = {
    'fixed_input': faf_image,
    'moving_input': slo_image
    }

    # compute registration
    A, _ = eyelinerp(data)

    # apply registration to image and mask
    reg_slo_image = eyelinerp.apply_transform(A.squeeze(0), slo_image.squeeze(0)) # (256, 256), [0, 1] 
    reg_slo_ga_seg = eyelinerp.apply_transform(A.squeeze(0), slo_ga_seg.squeeze(0)) # (256, 256), [0, 1]
 
    # convert to numpy
    faf_image = tensor2numpy(faf_image.cpu().squeeze(0) * 255).astype(np.uint8)
    faf_ga_seg = tensor2numpy(faf_ga_seg.cpu().squeeze(0) * 255).astype(np.uint8)
    slo_image = tensor2numpy(slo_image.cpu().squeeze(0) * 255).astype(np.uint8)
    slo_ga_seg = tensor2numpy(slo_ga_seg.cpu().squeeze(0) * 255).astype(np.uint8)
    reg_slo_image = tensor2numpy(reg_slo_image.cpu() * 255).astype(np.uint8)
    reg_slo_ga_seg = tensor2numpy(reg_slo_ga_seg.cpu() * 255).astype(np.uint8)

    # resize registered image to the original size of fixed image
    faf_image = cv2.resize(faf_image, (h, w), interpolation=cv2.INTER_CUBIC) # (1536, 1536), [0, 1]
    faf_ga_seg = cv2.resize(faf_ga_seg, (h, w), interpolation=cv2.INTER_CUBIC) # (1536, 1536), [0, 1]
    slo_image = cv2.resize(slo_image, (h, w), interpolation=cv2.INTER_CUBIC) # (1536, 1536), [0, 1]
    slo_ga_seg = cv2.resize(slo_ga_seg, (h, w), interpolation=cv2.INTER_CUBIC) # (1536, 1536), [0, 1]
    reg_slo_image = cv2.resize(reg_slo_image, (h, w), interpolation=cv2.INTER_CUBIC) # (1536, 1536), [0, 1]
    reg_slo_ga_seg = cv2.resize(reg_slo_ga_seg, (h, w), interpolation=cv2.INTER_CUBIC) # (1536, 1536), [0, 1]

    return faf_image, reg_slo_image, np.where(faf_ga_seg > 0, 255, 0).astype(np.uint8), np.where(reg_slo_ga_seg > 0, 255, 0).astype(np.uint8)

def convert_numpy_to_tempfile(image_array, save_as='file.png'):
    """
    Saves a NumPy array as a temporary image file in the specified format.
    Returns the path to the temporary file.
    """
    image = Image.fromarray(cv2.cvtColor(image_array, cv2.COLOR_BGR2RGB))
    image.save(save_as)
    return None

def draw_contours_on_image(image_array, segmentation_array):
    """
    Draws contours from the segmentation on the image and returns the resulting image as a NumPy array.
    """
    image = image_array.copy()
    contours, _ = cv2.findContours(segmentation_array, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    cv2.drawContours(image, contours, -1, (0, 255, 0), thickness=6)
    return image

def process_slide(data):
    """
    Process a single slide's data: generate converted images and overlay images.
    """
    image_1_path, image_2_path, image_1_seg, image_2_seg = data[:4]

    # register images
    image1, image2, seg1, seg2 = register(image_1_path, image_2_path, image_1_seg, image_2_seg)

    # Generate overlay images
    image_1_overlay = draw_contours_on_image(image1, seg1)
    image_2_overlay = draw_contours_on_image(image2, seg2)

    # Convert images to temporary files for PowerPoint
    convert_numpy_to_tempfile(image1, save_as='image_1_temp.png')
    convert_numpy_to_tempfile(seg1, save_as='image_1_seg_temp.png')
    convert_numpy_to_tempfile(image_1_overlay, save_as='image_1_overlay_temp.png')

    convert_numpy_to_tempfile(image2, save_as='image_2_temp.png')
    convert_numpy_to_tempfile(seg2, save_as='image_2_seg_temp.png')
    convert_numpy_to_tempfile(image_2_overlay, save_as='image_2_overlay_temp.png')

    return

def create_ppt(slides_data, ppt_output_path):
    prs = Presentation()

    for data in tqdm(slides_data):

        # get text data
        patient_id, laterality, exam_date = data[4:]

        # process images for the slide
        process_slide(data)

        # Add a slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = f"Patient ID: {patient_id}, Laterality: {laterality}, Exam Date: {exam_date}"

        # Add images to the slide
        top_row_y = Inches(2)
        bottom_row_y = Inches(4.5)

        slide.shapes.add_picture('image_1_temp.png', Inches(1.5), top_row_y, width=Inches(2))
        slide.shapes.add_picture('image_1_seg_temp.png', Inches(4), top_row_y, width=Inches(2))
        slide.shapes.add_picture('image_1_overlay_temp.png', Inches(6.5), top_row_y, width=Inches(2))

        slide.shapes.add_picture('image_2_temp.png', Inches(1.5), bottom_row_y, width=Inches(2))
        slide.shapes.add_picture('image_2_seg_temp.png', Inches(4), bottom_row_y, width=Inches(2))
        slide.shapes.add_picture('image_2_overlay_temp.png', Inches(6.5), bottom_row_y, width=Inches(2))

        # Save the presentation
        prs.save(ppt_output_path)

if __name__ == '__main__':

    # load data
    df = pd.read_csv('data/GA_progression_modelling_data_redone/clean_data_talisa_ga_10312024.csv')

    # get images
    slides_data = []
    for pat, pat_df in df.groupby('PID'):
        for lat, lat_df in pat_df.groupby('Laterality'):
            for date, date_df in lat_df.groupby('ExamDate'):
                df_af = date_df[date_df.Procedure == 'Af']
                df_ir = date_df[date_df['type'] == 'SLOImage']

                for i, row_af in df_af.iterrows():
                    for j, row_ir in df_ir.iterrows():
                        slides_data.append([row_af.file_path_coris, row_ir.file_path_coris, row_af.file_path_ga_seg, row_ir.file_path_ga_seg, pat, lat, date])

    # Create the PowerPoint presentation
    create_ppt(slides_data, "patient_slides.pptx")

    # remove temporary files
    os.remove('image_1_temp.png')
    os.remove('image_1_seg_temp.png')
    os.remove('image_1_overlay_temp.png')
    os.remove('image_2_temp.png')
    os.remove('image_2_seg_temp.png')
    os.remove('image_2_overlay_temp.png')