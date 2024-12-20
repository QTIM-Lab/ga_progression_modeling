import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from PIL import Image
import cv2
import os
from pptx import Presentation
from pptx.util import Inches
import io
import argparse
from concurrent.futures import ThreadPoolExecutor, ProcessPoolExecutor
from functools import partial
import warnings
warnings.filterwarnings('ignore')

def parse_arguments():
    """Parse command line arguments."""
    parser = argparse.ArgumentParser(description='Generate PowerPoint comparison of model segmentations.')
    
    # Required arguments
    parser.add_argument('--csv_dir', type=str, required=True,
                        help='Directory containing CSV files with model outputs')
    parser.add_argument('--output_path', type=str, required=True,
                        help='Path where the PowerPoint file will be saved')
    
    # Optional arguments with defaults
    parser.add_argument('--image_col', type=str, default='file_path_coris',
                        help='Column name containing original image paths (default: file_path_coris)')
    parser.add_argument('--seg_col', type=str, default='file_path_ga_seg',
                        help='Column name containing segmentation paths (default: file_path_ga_seg)')
    parser.add_argument('--pid_col', type=str, default='PID',
                       help='Column name for Patient ID (default: PID)')
    parser.add_argument('--lat_col', type=str, default='Laterality',
                       help='Column name for Laterality (default: Laterality)')
    parser.add_argument('--date_col', type=str, default='ExamDate',
                       help='Column name for Exam Date (default: ExamDate)')
    parser.add_argument('--contour_color', type=str, default='blue',
                        help='Color of segmentation contours (default: blue)')
    parser.add_argument('--contour_thickness', type=int, default=5,
                        help='Thickness of segmentation contours (default: 2)')
    parser.add_argument('--fig_width', type=float, default=5.0,
                        help='Width of each subplot in inches (default: 5.0)')
    parser.add_argument('--fig_height', type=float, default=5.0,
                        help='Height of each subplot in inches (default: 5.0)')
    
    # New arguments for testing
    parser.add_argument('--test_mode', action='store_true',
                        help='Run in test mode with subset of data')
    parser.add_argument('--start_row', type=int, default=0,
                        help='Starting row index (default: 0)')
    parser.add_argument('--num_rows', type=int, default=5,
                        help='Number of rows to process in test mode (default: 5)')
    parser.add_argument('--random_sample', action='store_true',
                        help='Randomly sample rows instead of sequential selection')
    parser.add_argument('--n_workers', type=int, default=4,
                        help='Number of worker processes (default: 4)')
    parser.add_argument('--batch_size', type=int, default=10,
                        help='Number of images to process in parallel (default: 10)')
    
    return parser.parse_args()

def load_and_preprocess_image(image_path, size=(1536, 1536)):
    """Load and preprocess image."""
    img = Image.open(image_path.replace('~', '/home/veturiy'))
    img = img.resize(size)
    return np.array(img)

def load_segmentation(seg_path, size=(1536, 1536)):
    """Load segmentation mask."""
    seg = Image.open(seg_path.replace('~', '/home/veturiy'))
    seg = seg.resize(size)
    return np.array(seg)

def get_contour_color(color_name):
    """Convert color name to RGB tuple."""
    color_map = {
        'blue': (255, 0, 0),
        'red': (0, 0, 255),
        'green': (0, 255, 0),
        'yellow': (0, 255, 255),
        'magenta': (255, 0, 255),
        'cyan': (255, 255, 0),
        'white': (255, 255, 255),
    }
    return color_map.get(color_name.lower(), (255, 0, 0))  # Default to blue if color not found

def create_comparison_figure(original_path, segmentation_paths, model_names, args):
    """Create a figure comparing original image with different segmentations."""
    original = load_and_preprocess_image(original_path)
    if original is None:
        print(f"Warning: Could not load original image: {original_path}")
        return None
    
    n_models = len(segmentation_paths)
    fig, axes = plt.subplots(1, n_models + 1, 
                            figsize=(args.fig_width*(n_models + 1), args.fig_height))
    
    # Handle case where there's only one comparison (axes not in array)
    if n_models == 0:
        axes = [axes]
    
    # Display original image
    axes[0].imshow(original)
    axes[0].set_title('Original Image')
    axes[0].axis('off')
    
    contour_color = get_contour_color(args.contour_color)
    
    # Display segmentations
    for idx, (seg_path, model_name) in enumerate(zip(segmentation_paths, model_names), 1):
        seg = load_segmentation(seg_path)
        if seg is not None:
            contours, _ = cv2.findContours(seg, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            
            img_with_contour = original.copy()
            cv2.drawContours(img_with_contour, contours, -1, contour_color, args.contour_thickness)
            
            axes[idx].imshow(img_with_contour)
            axes[idx].set_title(f'{model_name} Segmentation')
            axes[idx].axis('off')
        else:
            print(f"Warning: Could not load segmentation: {seg_path}")
    
    plt.tight_layout()
    return fig

def save_figure_to_pptx(prs, figure, info_text):
    """Save matplotlib figure to PowerPoint slide."""
    buf = io.BytesIO()
    figure.savefig(buf, format='png', bbox_inches='tight', dpi=300)
    buf.seek(0)
    
    # Use title-only layout
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title at top
    title = slide.shapes.add_textbox(
        left=Inches(0.5),
        top=Inches(0.5),
        width=prs.slide_width - Inches(1),
        height=Inches(0.5)
    )
    title.text = info_text
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center the title
    
    # Add picture below title
    top_margin = Inches(1.2)  # Space for title
    available_height = prs.slide_height - top_margin
    
    # Add picture and scale maintaining aspect ratio
    pic = slide.shapes.add_picture(buf, 0, top_margin, width=prs.slide_width)
    aspect_ratio = pic.height / pic.width
    new_width = prs.slide_width
    new_height = min(int(new_width * aspect_ratio), available_height)
    
    # Update size and center horizontally
    pic.width = new_width
    pic.height = new_height
    pic.left = 0  # Align to left edge
    pic.top = top_margin  # Position below title

def process_single_comparison(args_tuple):
    """Process a single comparison (for parallel processing)."""
    original_path, segmentation_paths, model_names, args = args_tuple
    
    # Create figure
    fig = create_comparison_figure(original_path, segmentation_paths, model_names, args)
    
    if fig is not None:
        # Save to bytes buffer
        buf = io.BytesIO()
        fig.savefig(buf, format='png', bbox_inches='tight', dpi=300)
        buf.seek(0)
        plt.close(fig)
        return buf
    return None

def main():
    args = parse_arguments()
    
    print("Loading CSV files...")
    # Load all CSVs at once and store in memory
    # Get the list of files in the directory along with their creation times
    files_with_times = [
        (file, os.path.getctime(os.path.join(args.csv_dir, file)))
        for file in os.listdir(args.csv_dir)
        if os.path.isfile(os.path.join(args.csv_dir, file)) and file.endswith('.csv')
    ]
    # Sort the files by creation time (ascending order)
    sorted_files = sorted(files_with_times, key=lambda x: x[1])
    csv_files = [file for file, _ in sorted_files]
    # csv_files = [f for f in os.listdir(args.csv_dir) if f.endswith('.csv')]
    if not csv_files:
        raise ValueError(f"No CSV files found in directory: {args.csv_dir}")
    
    model_names = [os.path.splitext(f)[0].split('af_')[-1] for f in csv_files]
    
    # Load all DataFrames into a dictionary
    dfs = {}
    for csv_file in csv_files:
        df = pd.read_csv(os.path.join(args.csv_dir, csv_file))
        dfs[csv_file] = df.set_index(args.image_col)  # Index by image path for faster lookup
    
    # Get first DataFrame
    first_df = dfs[csv_files[0]].reset_index()
    
    # Select rows based on test mode
    if args.test_mode:
        if args.random_sample:
            selected_rows = np.random.choice(
                first_df.index, 
                size=min(args.num_rows, len(first_df)), 
                replace=False
            )
            first_df = first_df.loc[selected_rows]
        else:
            first_df = first_df.iloc[args.start_row:args.start_row + args.num_rows]
    
    print(f"Processing {len(first_df)} images...")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Prepare arguments for parallel processing
    process_args = []
    for _, row in first_df.iterrows():
        original_path = row[args.image_col]
        
        # Get segmentation paths
        segmentation_paths = []
        for csv_file in csv_files:
            try:
                matching_row = dfs[csv_file].loc[original_path]
                segmentation_paths.append(matching_row[args.seg_col])
            except KeyError:
                print(f"Warning: No matching row found in {csv_file} for image {original_path}")
                continue
        
        process_args.append((original_path, segmentation_paths, model_names, args))
    
    # Process images in parallel batches
    with ProcessPoolExecutor(max_workers=args.n_workers) as executor:
        for i in range(0, len(process_args), args.batch_size):
            batch = process_args[i:i + args.batch_size]
            
            # Process batch in parallel
            results = list(executor.map(process_single_comparison, batch))
            
            # Create slides for successful results
            for j, buf in enumerate(results):
                if buf is not None:
                    idx = i + j
                    row = first_df.iloc[idx]
                    
                    # Create info text
                    info_text = (f"Patient ID: {row[args.pid_col]} | "
                               f"Laterality: {row[args.lat_col]} | "
                               f"Exam Date: {row[args.date_col]}")
                    
                    # Add to presentation
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # Add title
                    title_shape = slide.shapes.title
                    title_shape.text = info_text
                    
                    # Add and position image
                    pic = slide.shapes.add_picture(buf, 0, Inches(1), width=prs.slide_width)
                    aspect_ratio = pic.height / pic.width
                    new_width = prs.slide_width
                    new_height = int(new_width * aspect_ratio)
                    pic.width = new_width
                    pic.height = new_height
                    pic.left = int((prs.slide_width - pic.width) / 2)
                    remaining_height = prs.slide_height - title_shape.height
                    pic.top = title_shape.height + int((remaining_height - pic.height) / 2)
            
            print(f"Processed {min(i + args.batch_size, len(process_args))}/{len(process_args)} images")
    
    print("Saving presentation...")
    prs.save(args.output_path)
    print(f"Presentation saved to: {args.output_path}")

if __name__ == "__main__":
    main()

'''
# def main():
#     # Parse arguments
#     args = parse_arguments()
    
#     # Ensure output directory exists
#     output_dir = os.path.dirname(args.output_path)
#     if output_dir and not os.path.exists(output_dir):
#         os.makedirs(output_dir)
    
#     # Get all CSV files
#     csv_files = [f for f in os.listdir(args.csv_dir) if f.endswith('.csv')]
#     if not csv_files:
#         raise ValueError(f"No CSV files found in directory: {args.csv_dir}")
    
#     model_names = [os.path.splitext(f)[0].split('af_')[-1] for f in csv_files]
    
#     # Create presentation
#     prs = Presentation()
#     prs.slide_width = Inches(16)
#     prs.slide_height = Inches(9)
    
#     # Read first CSV to get list of images
#     first_df = pd.read_csv(os.path.join(args.csv_dir, csv_files[0]))
    
#     if args.image_col not in first_df.columns:
#         raise ValueError(f"Image column '{args.image_col}' not found in CSV")
#     if args.seg_col not in first_df.columns:
#         raise ValueError(f"Segmentation column '{args.seg_col}' not found in CSV")
    
#     # Select rows based on test mode settings
#     if args.test_mode:
#         if args.random_sample:
#             # Random sampling
#             total_rows = len(first_df)
#             if args.num_rows > total_rows:
#                 print(f"Warning: Requested {args.num_rows} rows but CSV only has {total_rows} rows")
#                 selected_rows = first_df.index
#             else:
#                 selected_rows = np.random.choice(first_df.index, 
#                                                size=min(args.num_rows, total_rows), 
#                                                replace=False)
#             first_df = first_df.loc[selected_rows]
#         else:
#             # Sequential selection
#             end_row = args.start_row + args.num_rows
#             first_df = first_df.iloc[args.start_row:end_row]
        
#         print(f"Running in test mode with {len(first_df)} images")
    
#     # Process each image
#     for idx, row in first_df.iterrows():
#         print(f"Processing image {idx + 1}/{len(first_df)}")
#         original_path = row[args.image_col]

#         # Create patient info text
#         info_text = (f"Patient ID: {row[args.pid_col]} | "
#                     f"Laterality: {row[args.lat_col]} | "
#                     f"Exam Date: {row[args.date_col]}")
        
#         # Collect segmentation paths from all models
#         segmentation_paths = []
#         for csv_file in csv_files:
#             df = pd.read_csv(os.path.join(args.csv_dir, csv_file))
#             matching_rows = df[df[args.image_col] == original_path]
            
#             if len(matching_rows) == 0:
#                 print(f"Warning: No matching row found in {csv_file} for image {original_path}")
#                 continue
                
#             segmentation_paths.append(matching_rows.iloc[0][args.seg_col])
        
#         # Create comparison figure
#         fig = create_comparison_figure(original_path, segmentation_paths, model_names, args)
        
#         if fig is not None:
#             save_figure_to_pptx(prs, fig, info_text)
#             plt.close(fig)
    
#     # Save presentation
#     prs.save(args.output_path)
#     print(f"\nPresentation saved to: {args.output_path}")

'''